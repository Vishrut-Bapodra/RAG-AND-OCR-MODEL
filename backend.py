# =============================
# Imports
# =============================
import os
import io
import re
import uuid
import json
from pathlib import Path
from typing import List, Dict, Any, Optional
import requests
import PyPDF2
import docx
from pptx import Presentation
import fitz  # PyMuPDF

# =============================
# Configuration
# =============================
BASE_DIR = Path(__file__).parent.resolve()
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

EMBED_MODEL_NAME = "intfloat/multilingual-e5-base"
CHUNK_MAX_CHARS = 900
CHUNK_OVERLAP_SENTENCES = 1

# =============================
# OCR & Text Extraction
# =============================

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Strategy:
    1. Try digital text extraction via PyPDF2
    2. If text is insufficient → use PyMuPDF OCR
    """
    # --- Try digital text ---
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        digital_text = "\n".join([p.extract_text() or "" for p in reader.pages]).strip()
        if len(digital_text) > 500:
            return digital_text
    except Exception:
        pass

    # --- OCR fallback (PyMuPDF full-page render) ---
    text_parts = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_index, page in enumerate(doc):
        pix = page.get_pixmap(dpi=300)
        img_bytes = pix.tobytes("png")

        try:
            import pytesseract
            from PIL import Image

            img = Image.open(io.BytesIO(img_bytes))
            page_text = pytesseract.image_to_string(img, lang="eng+guj")
            if page_text.strip():
                text_parts.append(page_text)
        except Exception as e:
            print(f"OCR failed on page {page_index}: {e}")

    return "\n".join(text_parts).strip()


def extract_text_from_docx(b: bytes) -> str:
    try:
        with io.BytesIO(b) as bio:
            doc = docx.Document(bio)
            return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def extract_text_from_pptx(b: bytes) -> str:
    try:
        texts = []
        with io.BytesIO(b) as bio:
            prs = Presentation(bio)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""


def extract_text(filename: str, file_bytes: bytes) -> str:
    fn = filename.lower()
    if fn.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    if fn.endswith(".docx"):
        return extract_text_from_docx(file_bytes)
    if fn.endswith(".pptx"):
        return extract_text_from_pptx(file_bytes)
    if fn.endswith(".txt"):
        return file_bytes.decode("utf-8", errors="ignore")
    return ""


# =============================
# Chunking (Sentence-based)
# =============================

_SENTENCE_SPLIT = re.compile(r"(?<=[.!?।])\s+")

def chunk_text(text: str) -> List[str]:
    sentences = [s.strip() for s in _SENTENCE_SPLIT.split(text) if s.strip()]
    chunks = []
    current = []

    for sent in sentences:
        current.append(sent)
        if sum(len(s) for s in current) >= CHUNK_MAX_CHARS:
            chunks.append(" ".join(current))
            current = current[-CHUNK_OVERLAP_SENTENCES:]

    if current:
        chunks.append(" ".join(current))

    return chunks


# =============================
# Vector Store (Lazy)
# =============================

class LazyVectorStore:
    def __init__(self):
        self.model = None
        self.index = None
        self.vectors = None
        self.metadatas = {}
        self.next_id = 0

    def _init(self):
        if self.model:
            return

        import numpy as np
        import faiss
        from sentence_transformers import SentenceTransformer

        self.model = SentenceTransformer(EMBED_MODEL_NAME)
        dim = self.model.get_sentence_embedding_dimension()
        self.index = faiss.IndexFlatIP(dim)
        self.vectors = np.zeros((0, dim), dtype="float32")

    def add(self, texts: List[str], metas: List[Dict[str, Any]]):
        if not texts:
            return

        self._init()
        import numpy as np

        emb = self.model.encode(texts, convert_to_numpy=True)
        emb = emb / (np.linalg.norm(emb, axis=1, keepdims=True) + 1e-9)
        emb = emb.astype("float32")

        self.index.add(emb)

        for i, meta in enumerate(metas):
            self.metadatas[self.next_id + i] = meta

        self.next_id += len(texts)

    def search(self, query: str, k: int = 5):
        self._init()
        import numpy as np

        q = self.model.encode([query], convert_to_numpy=True)
        q = q / (np.linalg.norm(q) + 1e-9)
        D, I = self.index.search(q.astype("float32"), k)

        results = []
        for score, idx in zip(D[0], I[0]):
            if idx == -1:
                continue
            meta = self.metadatas.get(int(idx), {})
            results.append({
                "score": float(score),
                "source": meta.get("source"),
                "chunk_index": meta.get("chunk_index"),
                "text": meta.get("text"),
            })
        return results


STORE = LazyVectorStore()

# =============================
# Indexing
# =============================

def save_and_index_file(uploaded_file):
    raw = uploaded_file.getvalue()
    stored_name = f"{uuid.uuid4().hex}_{uploaded_file.name}"
    (UPLOAD_DIR / stored_name).write_bytes(raw)

    text = extract_text(uploaded_file.name, raw)
    chunks = chunk_text(text)

    metas = [
        {
            "source": stored_name,
            "chunk_index": i,
            "text": c
        }
        for i, c in enumerate(chunks)
    ]

    STORE.add(chunks, metas)

    return {
        "file": uploaded_file.name,
        "status": "ok",
        "chunks_added": len(chunks),
    }


# =============================
# OpenRouter LLM Call
# =============================

def call_openrouter(
    api_key: str,
    model: str,
    messages: list,
    temperature: float = 0.3,
    max_tokens: int = 1024,
) -> str:
    url = "https://openrouter.ai/api/v1/chat/completions"

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost",
        "X-Title": "UE.ai RAG Bot",
    }

    payload = {
        "model": model,              
        "messages": messages,         
        "temperature": temperature,
        "max_tokens": max_tokens,
    }

    resp = requests.post(url, headers=headers, json=payload, timeout=60)
    resp.raise_for_status()

    data = resp.json()
    return data["choices"][0]["message"]["content"]




# =============================
# Main RAG Function
# =============================

def get_chatbot_response(
    messages: List[Dict[str, str]],
    top_k: int = 5,
    temperature: float = 0.3,
    synthesize: bool = True,
    openrouter_model: str = "openrouter/free",
):

    api_key = os.getenv("OPENROUTER_API_KEY")
    if not api_key:
        return "LLM API key not configured."

    user_query = next(
        (m["content"] for m in reversed(messages) if m["role"] == "user"),
        ""
    ).strip()

    if not user_query:
        return "Please ask a valid question."

    results = STORE.search(user_query, top_k)
    if not results:
        return "No relevant information found in the uploaded documents."

    context_blocks = []
    for i, r in enumerate(results, 1):
        context_blocks.append(
            f"[{i}] Source: {r['source']} (chunk {r['chunk_index']})\n{r['text']}"
        )

    context = "\n\n".join(context_blocks)

    prompt = [
        {
            "role": "system",
            "content": (
                "You are a legal and document analysis assistant. "
                "Answer clearly, structured, and factually using ONLY the provided context."
            ),
        },
        {
            "role": "user",
            "content": f"Question:\n{user_query}\n\nContext:\n{context}",
        },
    ]

    try:
        return call_openrouter(
            api_key=api_key,
            model=openrouter_model,
            messages=prompt,
            temperature=temperature,
        )
    except Exception as e:
        print("LLM failure:", e)
        return "An error occurred while generating the answer."
